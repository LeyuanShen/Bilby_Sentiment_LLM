from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PPT = BASE_DIR / "Alibaba_News_Sentiment_Presentation.pptx"
SCRIPT_TXT = BASE_DIR / "speech_script_2min.txt"
FIG_PATH = BASE_DIR / "output" / "sentiment_vs_returns.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "navy": RGBColor(19, 41, 75),
    "blue": RGBColor(50, 90, 168),
    "teal": RGBColor(34, 125, 140),
    "gold": RGBColor(201, 145, 0),
    "light": RGBColor(243, 246, 250),
    "dark": RGBColor(34, 34, 34),
    "muted": RGBColor(96, 103, 112),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(46, 125, 50),
}


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.bold = True
    run.font.size = Pt(25)
    run.font.color.rgb = COLORS["white"] if dark else COLORS["navy"]
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.02), Inches(9.5), Inches(0.5))
        p2 = sub_box.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = RGBColor(225, 231, 240) if dark else COLORS["muted"]


def add_bullets(slide, items, left, top, width, height, font_size=18, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS["dark"]
        p.space_after = Pt(12)
    return box


def add_metric_card(slide, left, top, width, height, value, label, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = "Aptos Display"
    r1.font.bold = True
    r1.font.size = Pt(24)
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Aptos"
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = COLORS["dark"]


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS["light"])
add_title(
    slide,
    "LLM-Based Sentiment Scoring for Alibaba News",
    "HKU Capstone Project | Turning financial news into a structured market signal",
)

banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.45), Inches(4.9), Inches(4.85))
banner.fill.solid()
banner.fill.fore_color.rgb = COLORS["navy"]
banner.line.color.rgb = COLORS["navy"]

quote = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(4.2), Inches(3.6))
qf = quote.text_frame
qf.word_wrap = True
p = qf.paragraphs[0]
r = p.add_run()
r.text = 'Goal: build a practical LLM system that reads Alibaba-related news and outputs a consistent 1-5 sentiment score.'
r.font.name = "Aptos Display"
r.font.bold = True
r.font.size = Pt(22)
r.font.color.rgb = COLORS["white"]

p2 = qf.add_paragraph()
p2.space_before = Pt(20)
r2 = p2.add_run()
r2.text = 'Why it matters: news moves expectations quickly, but manual labeling is slow, subjective, and hard to scale.'
r2.font.name = "Aptos"
r2.font.size = Pt(15)
r2.font.color.rgb = RGBColor(228, 235, 244)

add_bullets(
    slide,
    [
        "Input: Alibaba financial news title + body",
        "Output: sentiment score from 1 to 5",
        "Use case: transform unstructured text into a quantitative signal for market analysis",
    ],
    Inches(6.0), Inches(1.75), Inches(6.4), Inches(3.0), font_size=18, color=COLORS["dark"]
)

callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(5.2), Inches(6.2), Inches(0.95))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(232, 240, 252)
callout.line.color.rgb = COLORS["blue"]
ctf = callout.text_frame
cp = ctf.paragraphs[0]
cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
cr.text = "Core idea: use LLM prompting to produce scalable and consistent sentiment labels"
cr.font.name = "Aptos"
cr.font.bold = True
cr.font.size = Pt(16)
cr.font.color.rgb = COLORS["navy"]

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS["white"])
add_title(slide, "Methodology", "Few-shot prompting + financial validation")

steps = [
    ("1. Human-labeled data", "176 Alibaba news articles with sentiment labels"),
    ("2. Few-shot prompt", "Gemini 2.5 Flash Lite learns the scoring style from examples"),
    ("3. LLM scoring", "Each article is classified into a 1-5 sentiment score"),
    ("4. Validation", "Compare with human labels and test relation with stock excess returns"),
]

x_positions = [0.75, 3.1, 5.45, 7.8]
for idx, (title, desc) in enumerate(steps):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[idx]), Inches(2.0), Inches(2.0), Inches(2.25))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(243, 246, 250)
    card.line.color.rgb = COLORS["teal"] if idx % 2 else COLORS["blue"]
    tf = card.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Aptos Display"
    r1.font.bold = True
    r1.font.size = Pt(15)
    r1.font.color.rgb = COLORS["navy"]
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.name = "Aptos"
    r2.font.size = Pt(11)
    r2.font.color.rgb = COLORS["dark"]
    if idx < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x_positions[idx] + 2.05), Inches(2.68), Inches(0.4), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS["gold"]
        arrow.line.color.rgb = COLORS["gold"]

add_bullets(
    slide,
    [
        "Prompt persona: sell-side analyst covering Alibaba",
        "Score design: 1 = strongly negative, 5 = strongly positive",
        "Evaluation: classification accuracy and predictive value for T+1, T+3, T+5 excess returns",
    ],
    Inches(0.95), Inches(4.8), Inches(11.3), Inches(1.7), font_size=16
)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS["light"])
add_title(slide, "Results and Business Value", "Strong label alignment and useful market signal")

add_metric_card(slide, Inches(0.75), Inches(1.55), Inches(2.1), Inches(1.4), "86.1%", "Exact match on held-out test set", COLORS["blue"])
add_metric_card(slide, Inches(3.0), Inches(1.55), Inches(2.35), Inches(1.4), "69.0%", "Exact match on independent 200-article set", COLORS["teal"])
add_metric_card(slide, Inches(5.55), Inches(1.55), Inches(2.1), Inches(1.4), "100%", "Within +/-1 score on 200-article set", COLORS["green"])
add_metric_card(slide, Inches(7.85), Inches(1.55), Inches(2.1), Inches(1.4), "0.228", "Spearman correlation with T+1 excess return", COLORS["gold"])

if FIG_PATH.exists():
    slide.shapes.add_picture(str(FIG_PATH), Inches(0.8), Inches(3.15), width=Inches(7.4))

add_bullets(
    slide,
    [
        "The LLM is close to human judgment on sentiment classification.",
        "Scores also contain useful information about short-term market reaction.",
        "The full pipeline has been packaged for direct company use.",
    ],
    Inches(8.55), Inches(3.35), Inches(4.0), Inches(2.2), font_size=16
)

final_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(5.65), Inches(3.95), Inches(0.95))
final_box.fill.solid()
final_box.fill.fore_color.rgb = COLORS["navy"]
final_box.line.color.rgb = COLORS["navy"]
ft = final_box.text_frame
fp = ft.paragraphs[0]
fp.alignment = PP_ALIGN.CENTER
fr = fp.add_run()
fr.text = "Outcome: a scalable news-to-signal workflow for financial analysis"
fr.font.name = "Aptos"
fr.font.bold = True
fr.font.size = Pt(15)
fr.font.color.rgb = COLORS["white"]

prs.save(OUTPUT_PPT)

speech = """Good afternoon everyone.

Our project is about using a large language model to analyze the sentiment of Alibaba-related financial news.

We started from a simple problem. Financial news can affect market expectations very quickly, but reading and labeling news manually takes a lot of time. It is also hard to keep the scoring consistent when the amount of news becomes large.

So our goal was to build a system that can automatically read a news article and give it a sentiment score from 1 to 5.

For the method, we used Google Gemini 2.5 Flash Lite. We also designed a few-shot prompt based on manually labeled Alibaba news. In this way, the model could learn our scoring style and follow the same standard more closely.

For each article, the model reads the title and the body, and then outputs a structured sentiment score. After that, we evaluated the model in two ways. First, we compared its predictions with human labels. Second, we checked whether the sentiment scores were related to Alibaba's future stock performance.

Our results were quite encouraging. On the original test set, the model achieved 86.1 percent exact-match accuracy. On another independent dataset with 200 articles, it achieved 69.0 percent exact-match accuracy, and 100 percent accuracy within plus or minus one score.

We also found that the predicted sentiment score had a positive correlation with Alibaba's next-day excess return, with a value of 0.228. This suggests that the score is not only close to human judgment, but also contains useful market information.

In summary, our project shows that an LLM can turn unstructured financial news into a structured sentiment signal. It is faster than manual labeling, easier to scale, and can support future financial analysis and decision-making.

Thank you."""
SCRIPT_TXT.write_text(speech, encoding="utf-8")
print(f"Saved PPT to: {OUTPUT_PPT}")
print(f"Saved script to: {SCRIPT_TXT}")